'''
this .py implements the data preparation stage of the RAG pipeline described in
the dissertation methodology (data preparation and chunking). it reads the cleaned Word, PowerPoint
and PDF course files, strips the storyboard production apparatus that is not
course content, chunks what is left into ~500 word windows with 50 word overlap,
and tags every chunk with metadata (week, topic, artefact type, source file)

Output: knowledge_base_chunks.jsonl, one JSON object per line, ready to be
embedded into ChromaDB by build_chromadb.py

Usage:
  py -m pip install python-docx python-pptx pypdf
  py extract_and_chunk.py

v2 rewrite. the four changes are:

1. merged cell duplication.
   v1 read tables with "for row in table.rows: for cell in row.cells:". python-docx
   hands back the SAME cell object once per grid column it spans, so any text in a
   merged cell was written out once per span. the KEATS storyboards are built almost
   entirely from wide merged cells, so in 1.6 Network structure 68 distinct cells
   produced 148 visits and one cell was emitted 37 times. measured across the whole
   knowledge base, 58% of v1 was duplicated text and the five worst files were 74-82%
   duplicate. fixed by walking the document XML directly and visiting each <w:tc> once

2. document order.
   v1 wrote out every paragraph first and then every table, so table content was
   detached from the prose that introduced it and chunk boundaries landed in
   meaningless places. fixed with a single in order walk of the document body

3. production apparatus.
   storyboard reviewer instructions, video editing notes and F2F video links were
   being embedded as if they were course content (67, 12 and 26 occurrences). these
   are now dropped by pattern. alt text descriptions are deliberately KEPT and
   relabelled "Figure:", because they are human written descriptions of diagrams the
   extractor cannot see and are the only representation of those figures in the KB

4. speaker notes.
   v1 read PowerPoint shape text frames only, so any narration in the notes pane was
   lost. notes_slide is now read per slide
'''

import os
import re
import json
import argparse

from docx import Document
from docx.oxml.ns import qn
from pptx import Presentation
from pypdf import PdfReader

#paths come from config.py so the project folder stays portable.
#both can still be overridden at the command line with --source and --output
from config import CLEANED_DIR, CHUNKS_FILE

SOURCE_DIR = str(CLEANED_DIR)
OUTPUT_FILE = str(CHUNKS_FILE)

CHUNK_SIZE = 500      #words
CHUNK_OVERLAP = 50    #words

#only run the near duplicate guard on lines at least this long, so short repeated
#headings like "Notes" or "Task 1" are not wrongly thrown away
DEDUPE_MIN_WORDS = 12


#Page furniture and KEATS/Moodle quiz attempt chrome, not question content
BOILERPLATE_PATTERNS = [
    r"^Key links$",
    r"^Terms of Use$",
    r"^Privacy Policy$",
    r"^Contact Us$",
    r"^Student Data Collection Notice$",
    r"^Accessibility Statement$",
    r"^Release Notes$",
    r"^Social media$",
    r"King['\u2019]?s College London, Strand, London",
    r"^Tel \+44",
    r"^Jump to\.\.\.",
    r"^\u25c4",
    r"\u25ba\s*$",
    r"^\d{4} King['\u2019]?s College London",
    r"^Status\s*Finished",
    r"^Started\w*day",
    r"^Completed\w*day",
    r"^Duration\s*\d",
    r"^Grade\s*[\d.]+\s*out of",
    r"^Not answered$",
    r"^Marked out of \d",
    r"^Question\s*\d+$",
]

#Storyboard production apparatus. written for the media team, not for students
PRODUCTION_PATTERNS = [
    r"Instructions?\b.*\bfor storyboard reviewers",
    r"Changes made to the page itself will save automatically",
    r"please see the King['\u2019]?s style guide",
    r"^\[KEATS (book|page|activity)\b",
    r"switch from Editing to Reviewing mode",
    r"Highlight any text you wish to comment on",
    r"select New comment from the",
    r"We recommend opening and working with the document in the App",
    r"Editing notes from the ML",
    r"^Cut/remove from",
    r"^Continue [Vv]ideo from",
    r"^Video link in the F2F module",
    r"^Link to video:?$",
    r"^\[Embedded video\]",
    r"^Video details:?",
    r"^Link to video script:?",
    r"^Add details for thumbnail",
    r"^Awaiting final editing$",
    r"^\[?Download\]? file",
    r"kaf\.keats\.kcl\.ac\.uk",
    r"^\[?Instructional alert\]?$",
    r"^Caption \(optional\):?$",
    r"^Alt text:?$",
    r"^Slide \d+$",
    r"^TBC$",
    r"^N/?A$",
]

#compiled once rather than per line, since is_boilerplate() runs on every line of
#every file in the corpus
BOILER_RE = [re.compile(p, re.IGNORECASE) for p in BOILERPLATE_PATTERNS]
PROD_RE = [re.compile(p, re.IGNORECASE) for p in PRODUCTION_PATTERNS]

#markers stripped from inside a line that is otherwise worth keeping
INLINE_STRIP = [
    (re.compile(r"\[Instructional alert\]\s*", re.IGNORECASE), ""),
    (re.compile(r"^Alt text:\s*", re.IGNORECASE), "Figure: "),
    (re.compile(r"^Caption \(optional\):\s*", re.IGNORECASE), "Caption: "),
    (re.compile(r"https?://keats\.kcl\.ac\.uk/\S*", re.IGNORECASE), ""),
]


def clean_line(line: str) -> str:
    """strip the inline production markers but keep the readable content"""
    text = line.strip()
    for pattern, replacement in INLINE_STRIP:
        text = pattern.sub(replacement, text)
    return re.sub(r"\s+", " ", text).strip()


def is_boilerplate(line: str) -> bool:
    """True if the line is page furniture or storyboard production apparatus"""
    text = line.strip()
    if not text:
        return True
    for rx in BOILER_RE:
        if rx.search(text):
            return True
    for rx in PROD_RE:
        if rx.search(text):
            return True
    return False


def keep_line(text: str, kept: list, seen: set) -> None:
    """
    shared filter used by all three extractors. drops boilerplate and applies the
    near duplicate guard, so a repeated long line is only stored once per file
    """
    if not text or is_boilerplate(text):
        return
    words = text.split()
    if len(words) >= DEDUPE_MIN_WORDS:
        key = " ".join(w.lower() for w in words)
        if key in seen:
            return
        seen.add(key)
    kept.append(text)


def paragraph_text(p_el) -> str:
    """
    text of a single <w:p>, skipping any mc:Fallback subtree

    Word stores a text box twice, once under mc:Choice for modern readers and once
    under mc:Fallback for older ones, so collecting every <w:t> naively emits the
    text box content twice
    """
    parts = []
    #python-docx qn() has no 'mc' prefix registered, so the full URI is needed here
    fallback = "{http://schemas.openxmlformats.org/markup-compatibility/2006}Fallback"

    def walk(node):
        for child in node:
            if child.tag == fallback:
                continue
            if child.tag == qn("w:t"):
                parts.append(child.text or "")
            elif child.tag == qn("w:tab"):
                parts.append("\t")
            elif child.tag in (qn("w:br"), qn("w:cr")):
                parts.append(" ")
            else:
                walk(child)

    walk(p_el)
    return "".join(parts)


def walk_block(parent, lines: list) -> None:
    """
    recursively pull text out of a block level container in document order

    handles paragraphs and tables, descends into nested tables, and visits each
    <w:tc> exactly once. that last part is the fix for the merged cell duplication
    described at the top of this file, because reading direct <w:tc> children of a
    row gives each cell once, unlike python-docx row.cells which repeats a merged
    cell for every column it spans
    """
    for child in parent:
        if child.tag == qn("w:p"):
            lines.append(paragraph_text(child))
        elif child.tag == qn("w:tbl"):
            for row in child.findall(qn("w:tr")):
                for cell in row.findall(qn("w:tc")):
                    walk_block(cell, lines)
        elif child.tag == qn("w:sdt"):
            content = child.find(qn("w:sdtContent"))
            if content is not None:
                walk_block(content, lines)


def extract_docx_text(path: str) -> str:
    """extract a Word file in document order, one pass over the body"""
    doc = Document(path)
    raw_lines = []
    walk_block(doc.element.body, raw_lines)

    kept, seen = [], set()
    for raw in raw_lines:
        keep_line(clean_line(raw), kept, seen)
    return "\n".join(kept)


def extract_pptx_text(path: str) -> str:
    """extract a PowerPoint file including table cells and the speaker notes pane"""
    prs = Presentation(path)
    kept, seen = [], set()

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    keep_line(clean_line("".join(run.text for run in para.runs)), kept, seen)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    for cell in row.cells:
                        keep_line(clean_line(cell.text), kept, seen)

        #v1 dropped the notes pane completely
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame
            if notes is not None and notes.text.strip():
                for line in notes.text.split("\n"):
                    keep_line(clean_line(line), kept, seen)

    return "\n".join(kept)


def extract_pdf_text(path: str) -> str:
    """extract a PDF page by page, dropping the KEATS quiz attempt chrome"""
    reader = PdfReader(path)
    kept, seen = [], set()
    for page in reader.pages:
        for raw in (page.extract_text() or "").split("\n"):
            keep_line(clean_line(raw), kept, seen)
    return "\n".join(kept)


def classify_artefact_type(filename: str) -> str:
    """
    tag each file with what kind of teaching artefact it is, so the retrieval layer
    can filter on artefact_type as well as semantic similarity later on
    """
    name = filename.lower()
    if "knowledge_check" in name or "knowledge check" in name:
        return "quiz"
    if "quiz" in name:
        return "quiz"
    if "webinar" in name:
        return "webinar"
    if "virtual_lab" in name or "virtual lab" in name:
        return "lab_instructions"
    if "reflect" in name:
        return "reflection"
    if "discussion" in name:
        return "discussion"
    if "slides" in name:
        return "slides"
    if "outline" in name:
        return "outline"
    if "overview" in name or "learning_out" in name:
        return "overview"
    if re.search(r"\bexam\b", name):
        return "exam"
    return "lecture"


def clean_topic_label(filename: str) -> str:
    """turn a raw filename into something readable for the sources panel"""
    name = os.path.splitext(filename)[0]
    name = re.sub(r"^Cyber_M3_\d+(\.\d+)*_?", "", name)
    name = re.sub(r"^Cyber_MSc_Cyber_Security_", "", name)
    name = name.replace("_", " ").strip()
    return name if name else filename


def get_week_number(filepath: str, source_root: str) -> str:
    """
    work out which teaching week a file belongs to, first from the folder it sits
    in, then falling back to the filename. the four module wide Virtual Labs
    instruction files legitimately have no week and return "unknown"
    """
    rel = os.path.relpath(filepath, source_root)
    for part in rel.split(os.sep):
        match = re.match(r"Week\s*(\d+)", part, re.IGNORECASE)
        if match:
            return match.group(1)

    base = os.path.basename(filepath)
    #order matters here. the module numbering has to be checked before any "Week n"
    #text in the filename, otherwise a cross week artefact like
    #Cyber_M3_5_14_Quiz_week_1-5.docx gets tagged Week 1 off the words "week_1"
    #when it is really a Week 5 artefact covering weeks 1 to 5

    #e.g. Cyber_M3_1.6 or Cyber_M3_1_6
    match = re.search(r"M3[_.](\d+)[._]", base)
    if match:
        return match.group(1)
    #e.g. 3_16_Quiz_week_1-3 or 5_5_1_Quiz_5
    match = re.search(r"^(\d+)[_.]\d+", base)
    if match:
        return match.group(1)
    #e.g. Cyber_M3_Week_1_outline.docx
    match = re.search(r"Week[_\s]*(\d+)", base, re.IGNORECASE)
    if match:
        return match.group(1)
    return "unknown"


def chunk_text(text: str, chunk_size: int = CHUNK_SIZE, overlap: int = CHUNK_OVERLAP):
    """
    split into fixed size word windows with a small overlap, so a concept that
    straddles a boundary still appears whole in one of the two chunks

    fixed size chunking with overlap is kept deliberately. Qu, Tu and Bao (2025)
    found the extra computational cost of semantic chunking is not reliably repaid
    in retrieval quality, and Das et al. (2025) reported fixed windows with metadata
    filtering to be sufficient for small, well scoped corpora, which is exactly the
    regime of this five week knowledge base
    """
    words = text.split()
    if not words:
        return []
    chunks = []
    start = 0
    while start < len(words):
        end = start + chunk_size
        chunks.append(" ".join(words[start:end]))
        if end >= len(words):
            break
        start = end - overlap
    return chunks


def main():
    parser = argparse.ArgumentParser(
        description="Extract and chunk the Network Security knowledge base."
    )
    parser.add_argument("--source", default=SOURCE_DIR, help="folder of cleaned source files")
    parser.add_argument("--output", default=OUTPUT_FILE, help="output .jsonl path")
    args = parser.parse_args()

    if not os.path.isdir(args.source):
        print(f"ERROR: Source folder not found:\n  {args.source}")
        return

    all_chunks = []
    files_processed = 0
    files_skipped = 0

    for dirpath, dirnames, filenames in os.walk(args.source):
        for filename in sorted(filenames):
            filepath = os.path.join(dirpath, filename)
            ext = os.path.splitext(filename)[1].lower()

            if ext not in (".docx", ".pptx", ".pdf"):
                files_skipped += 1
                continue

            #"~$" files are Word lock files left behind by an open document
            if filename.strip().upper().startswith("DELETED") or filename.startswith("~$"):
                print(f"Excluded: {filename}")
                files_skipped += 1
                continue

            week = get_week_number(filepath, args.source)
            artefact_type = classify_artefact_type(filename)
            topic_label = clean_topic_label(filename)

            try:
                if ext == ".docx":
                    text = extract_docx_text(filepath)
                elif ext == ".pptx":
                    text = extract_pptx_text(filepath)
                else:
                    text = extract_pdf_text(filepath)
            except Exception as e:
                print(f"  ERROR reading {filename}: {e}")
                files_skipped += 1
                continue

            if not text.strip():
                print(f"  WARNING: no text extracted from {filename}")
                files_skipped += 1
                continue

            chunks = chunk_text(text)
            print(f"Processed: {filename} -> {len(chunks)} chunk(s) [Week {week}, {artefact_type}]")

            for i, chunk in enumerate(chunks):
                all_chunks.append({
                    "chunk_id": f"{os.path.splitext(filename)[0]}_chunk{i+1}",
                    "week": week,
                    "topic_label": topic_label,
                    "artefact_type": artefact_type,
                    "source_file": filename,
                    "chunk_index": i + 1,
                    "total_chunks_in_file": len(chunks),
                    "word_count": len(chunk.split()),
                    "text": chunk,
                })

            files_processed += 1

    with open(args.output, "w", encoding="utf-8") as f:
        for chunk in all_chunks:
            f.write(json.dumps(chunk, ensure_ascii=False) + "\n")

    print("\n" + "="*60)
    print("Done.")
    print(f"  Files processed: {files_processed}")
    print(f"  Files skipped:   {files_skipped}")
    print(f"  Total chunks:    {len(all_chunks)}")
    print(f"  Total words:     {sum(c['word_count'] for c in all_chunks)}")
    print(f"\nOutput saved to:\n  {args.output}")
    print("\nNext: rerun build_chromadb.py to rebuild the vector store.")


if __name__ == "__main__":
    main()


'''Sources for the extraction behaviour worked around in this file

python-docx (2025) Working with tables. [online]
Available at: https://python-docx.readthedocs.io/en/latest/user/tables.html
and https://python-docx.readthedocs.io/en/latest/dev/analysis/features/table/cell-merge.html
-Documents that row.cells returns the same cell object once per grid column a
merged cell spans. walk_block() reads the direct w:tc children of each w:tr
instead, so every cell is visited once.

python-openxml (2016) duplicate cells returned from row.cells. python-docx issue
#344. [online] Available at: https://github.com/python-openxml/python-docx/issues/344
-Confirms the duplication behaviour above is expected rather than a defect, which
is why it is handled here rather than reported upstream.

python-openxml (2015) how to read paragraphs AND tables? python-docx issue #276.
[online] Available at: https://github.com/python-openxml/python-docx/issues/276
-Establishes that python-docx exposes paragraphs and tables as separate
collections with no ordered view of the two. walk_block() recovers document order
by walking the body XML directly rather than by the iter_block_items helper
proposed in that thread.

Ecma International (2016) ECMA-376 Office Open XML File Formats, Part 3: Markup
Compatibility and Extensibility. 5th edn. [online] Available at:
https://ecma-international.org/publications-and-standards/standards/ecma-376/
-Specification for the mc:Choice and mc:Fallback pattern that stores text box
content twice. paragraph_text() skips the mc:Fallback subtree.
'''
